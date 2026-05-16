import networkx as nx
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

steps = [
    {"nodes": ["A"], "edges": []},
    {"nodes": ["B"], "edges": [("A", "B")]},
    {"nodes": ["C"], "edges": [("B", "C")]},
    {"nodes": ["D"], "edges": [("C", "D")]},
    {"nodes": ["E"], "edges": [("C", "E")]},
    {"nodes": ["F"], "edges": [("E", "F")]},
    {"nodes": ["G"], "edges": [("D", "G"), ("F", "G")]},
    {"nodes": ["H"], "edges": [("C", "H")]},
    {"nodes": ["I"], "edges": [("H", "I")]},
    {"nodes": ["J"], "edges": [("G", "J"), ("I", "J")]},
    {"nodes": ["K"], "edges": [("J", "K")]}
]

pos = {
    "A": (0, 5),
    "B": (1, 5),
    "C": (2, 4),
    "D": (3, 4),
    "E": (2, 3),
    "F": (3, 3),
    "G": (4, 4),
    "H": (2, 2),
    "I": (3, 2),
    "J": (5, 3.5),
    "K": (6, 3.5),
}

node_colors = {
    "A": "gray",
    "B": "gray",
    "C": "orange",
    "D": "orange",
    "E": "blue",
    "F": "blue",
    "G": "blue",
    "H": "red",
    "I": "red",
    "J": "green",  # release
    "K": "green"
}

branch_labels = {
    "B": "master",
    "D": "feature",
    "F": "dev",
    "H": "hotfix",
    "K": "release"
}

def compute_graph_metrics(G):
    num_vertices = G.number_of_nodes()
    num_edges = G.number_of_edges()
    in_degrees = dict(G.in_degree())
    out_degrees = dict(G.out_degree())
    merge_count = sum(1 for deg in in_degrees.values() if deg > 1)
    avg_degree = (sum(in_degrees.values()) + sum(out_degrees.values())) / num_vertices if num_vertices > 0 else 0
    in_deg = max(in_degrees.values()) if in_degrees else 0
    out_deg = max(out_degrees.values()) if out_degrees else 0
    depths = nx.single_source_shortest_path_length(G, source='A') if 'A' in G.nodes else {}
    min_depth = min(depths.values()) if depths else 0
    max_depth = max(depths.values()) if depths else 0
    return {
        "numberOfVertices": num_vertices,
        "numberOfEdges": num_edges,
        "inDegree": in_deg,
        "outDegree": out_deg,
        "averageDegree": round(avg_degree, 2),
        "mergeCount": merge_count,
        "minDepthOfCommitHistory": min_depth,
        "maxDepthOfCommitHistory": max_depth
    }

all_nodes = []
all_edges = []
img_paths = []

for i, step in enumerate(steps):
    all_nodes += step["nodes"]
    all_edges += step["edges"]

    G = nx.DiGraph()
    G.add_nodes_from(all_nodes)
    G.add_edges_from(all_edges)

    metrics = compute_graph_metrics(G)
    print(f"\nStep {i + 1}: Metrics for nodes {', '.join(step['nodes'])}")
    for key, val in metrics.items():
        print(f"  {key}: {val}")

    plt.figure(figsize=(12, 7))
    colors = [node_colors.get(n, "white") for n in G.nodes]

    nx.draw(
        G, pos,
        with_labels=True,
        node_color=colors,
        node_size=1000,
        font_weight='bold',
        arrows=True,
        arrowstyle='-|>',
        connectionstyle='arc3,rad=0.15',
        linewidths=1.5
    )

    for node, label in branch_labels.items():
        if node in G.nodes:
            x, y = pos[node]
            plt.text(x, y + 0.3, label, fontsize=10, ha='center', fontstyle='italic')

    plt.title(f"Step {i + 1}: Adding {', '.join(step['nodes'])}", fontsize=14)

    metrics_text = "\n".join([
        f"Vertices: {metrics['numberOfVertices']}",
        f"Edges: {metrics['numberOfEdges']}",
        f"InDegree: {metrics['inDegree']}",
        f"OutDegree: {metrics['outDegree']}",
        f"AvgDegree: {metrics['averageDegree']}",
        f"Merge Count: {metrics['mergeCount']}",
        f"Min Depth: {metrics['minDepthOfCommitHistory']}",
        f"Max Depth: {metrics['maxDepthOfCommitHistory']}"
    ])
    plt.gcf().text(0.75, 0.5, metrics_text, fontsize=10, bbox=dict(facecolor='white', alpha=0.8))

    plt.axis("off")
    img_path = f"ppt_git_step_{i + 1}.png"
    img_paths.append(img_path)
    plt.savefig(img_path, bbox_inches='tight')
    plt.close()

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]

for path in img_paths:
    slide = prs.slides.add_slide(blank_slide_layout)
    slide.shapes.add_picture(path, Inches(0.5), Inches(0.5), width=Inches(8.5))

output_pptx = "git_commit_graph_complex.pptx"
prs.save(output_pptx)

